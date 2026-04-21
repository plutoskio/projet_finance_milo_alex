from __future__ import annotations

import statistics
from pathlib import Path

from openpyxl import Workbook, load_workbook
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SOURCES = ROOT / "sources"
EXCEL_OUT = ROOT / "excel_model" / "Michelin_valuation_model.xlsx"
PPT_OUT = ROOT / "slides" / "Michelin_valuation_deck.pptx"
SUMMARY_OUT = ROOT / "exports" / "valuation_summary.md"

SOURCE_XLSX = SOURCES / "michelin_key_figures_2020_2025.xlsx"

BLUE = "1F4E5F"
TEAL = "2A9D8F"
RED = "D1495B"
GOLD = "F4A261"
LIGHT = "F6F8FA"
MID = "E6EBEF"
DARK = "20252B"
GRAY = "6B7280"


def pct(x: float, digits: int = 1) -> str:
    return f"{x * 100:.{digits}f}%"


def eur_bn(x: float, digits: int = 1) -> str:
    return f"EUR {x / 1000:.{digits}f} Md"


def eur(x: float, digits: int = 1) -> str:
    return f"EUR {x:.{digits}f}"


def read_michelin_history() -> dict:
    wb = load_workbook(SOURCE_XLSX, data_only=True)
    ws = wb["Key figures"]
    years = [2020, 2021, 2022, 2023, 2024, 2025]

    rows = {}
    header = [cell.value for cell in ws[1]]
    year_to_col = {int(v): idx for idx, v in enumerate(header) if isinstance(v, int)}

    wanted = {
        "Sales": "sales",
        "% change": "sales_growth",
        "Segment EBITDA (1)": "ebitda",
        "Segments operating income": "segment_ebit",
        "Segments operating margin": "segment_margin",
        "Operating income": "operating_income",
        "Net income": "net_income",
        "Income tax": "income_tax",
        "Effective tax rate": "tax_rate",
        "Cash flows from operating activities": "cfo",
        "Gross purchases of intangible assets and PP&E": "capex",
        "Net debt (2)": "net_debt",
        "Free cash flow (4)": "fcf",
        "Research and development expenses": "rd",
        "Number of employees (full time equivalent)": "employees",
    }

    for row in ws.iter_rows(values_only=True):
        label = row[0]
        if label in wanted:
            key = wanted[label]
            rows[key] = {year: row[year_to_col[year]] for year in years}

    market_ws = wb["Markets"]
    markets = {
        "PLT OE": {},
        "PLT Replacement": {},
        "Truck OE": {},
        "Truck Replacement": {},
        "Mining": {},
        "Agricultural": {},
    }
    blocks = {
        "PLT OE": range(4, 10),
        "PLT Replacement": range(19, 25),
        "Truck OE": range(34, 40),
        "Truck Replacement": range(49, 55),
    }
    for name, rows_range in blocks.items():
        for r in rows_range:
            year = market_ws.cell(r, 1).value
            if isinstance(year, int):
                values = [market_ws.cell(r, c).value or 0 for c in range(2, 7)]
                markets[name][year] = sum(values)

    for r in range(63, 70):
        year = market_ws.cell(r, 1).value
        if isinstance(year, int):
            markets["Mining"][year] = market_ws.cell(r, 2).value
    for r in range(75, 82):
        year = market_ws.cell(r, 1).value
        if isinstance(year, int):
            markets["Agricultural"][year] = market_ws.cell(r, 2).value

    segments_ws = wb["Reporting segments"]
    segments = {
        "Automobile et distribution associee": {},
        "Transport routier et distribution associee": {},
        "Specialites et distribution associee": {},
    }
    segment_ranges = {
        "Automobile et distribution associee": range(4, 9),
        "Transport routier et distribution associee": range(14, 19),
        "Specialites et distribution associee": range(24, 29),
    }
    for name, rows_range in segment_ranges.items():
        for r in rows_range:
            year = segments_ws.cell(r, 1).value
            if isinstance(year, int):
                segments[name][year] = {
                    "segment_ebit": segments_ws.cell(r, 2).value,
                    "margin": segments_ws.cell(r, 3).value,
                }

    return {"years": years, "rows": rows, "markets": markets, "segments": segments}


def build_calculations(history: dict) -> dict:
    rows = history["rows"]
    sales_2025 = rows["sales"][2025]
    ebit_2025 = rows["segment_ebit"][2025]
    ebitda_2025 = rows["ebitda"][2025]
    net_debt_2025 = rows["net_debt"][2025]
    shares_m = 684.74
    current_share_price = 32.24
    tax_rate = rows["tax_rate"][2025]

    forecast_years = [2026, 2027, 2028, 2029, 2030]
    growth = {2026: 0.015, 2027: 0.025, 2028: 0.030, 2029: 0.030, 2030: 0.025}
    ebit_margin = {2026: 0.107, 2027: 0.112, 2028: 0.116, 2029: 0.119, 2030: 0.120}
    da_margin = {2026: 0.073, 2027: 0.072, 2028: 0.071, 2029: 0.070, 2030: 0.070}
    capex_margin = {2026: 0.073, 2027: 0.072, 2028: 0.071, 2029: 0.070, 2030: 0.070}
    nwc_pct_incremental_sales = 0.010

    forecast = {}
    prev_sales = sales_2025
    for year in forecast_years:
        sales = prev_sales * (1 + growth[year])
        ebit = sales * ebit_margin[year]
        da = sales * da_margin[year]
        ebitda = ebit + da
        nopat = ebit * (1 - tax_rate)
        capex = sales * capex_margin[year]
        delta_nwc = (sales - prev_sales) * nwc_pct_incremental_sales
        fcff = nopat + da - capex - delta_nwc
        forecast[year] = {
            "sales": sales,
            "growth": growth[year],
            "ebit_margin": ebit_margin[year],
            "ebit": ebit,
            "da_margin": da_margin[year],
            "da": da,
            "ebitda": ebitda,
            "ebitda_margin": ebitda / sales,
            "tax_rate": tax_rate,
            "nopat": nopat,
            "capex_margin": capex_margin[year],
            "capex": capex,
            "delta_nwc": delta_nwc,
            "fcff": fcff,
        }
        prev_sales = sales

    comps = [
        {
            "company": "Bridgestone",
            "ticker": "TYO:5108",
            "country": "Japon",
            "activity": "Pneu mondial, tourisme / truck / solutions",
            "relevance": "Leader mondial pneus, exposition replacement/OE comparable",
            "limitation": "Exposition Japon / Asie plus marquee",
            "market_cap": 4320,
            "enterprise_value": 4490,
            "ev_eqv": 4490 / 4320,
            "ev_sales": 1.01,
            "ev_ebitda": 5.38,
            "ev_ebit": 9.34,
            "ebit_margin": 0.1087,
        },
        {
            "company": "Goodyear",
            "ticker": "NASDAQ:GT",
            "country": "Etats-Unis",
            "activity": "Pneumatiques purs, forte exposition US",
            "relevance": "Pneumatiques purs, mais levier et marge faible",
            "limitation": "Levier eleve et marge tres faible, donc multiples parfois deformes",
            "market_cap": 2060,
            "enterprise_value": 8480,
            "ev_eqv": 8480 / 2060,
            "ev_sales": 0.46,
            "ev_ebitda": 6.80,
            "ev_ebit": 23.48,
            "ebit_margin": 0.0197,
        },
        {
            "company": "Continental",
            "ticker": "ETR:CON",
            "country": "Allemagne",
            "activity": "Equipementier auto diversifie avec activite pneus",
            "relevance": "Equipementier auto diversifie avec activite pneus",
            "limitation": "Groupe moins pur pneu que Michelin",
            "market_cap": 13540,
            "enterprise_value": 19080,
            "ev_eqv": 19080 / 13540,
            "ev_sales": 0.97,
            "ev_ebitda": 6.89,
            "ev_ebit": 12.18,
            "ebit_margin": 0.0792,
        },
        {
            "company": "Pirelli",
            "ticker": "BIT:PIRC",
            "country": "Italie",
            "activity": "Pneu premium et performance",
            "relevance": "Pneus premium, forte proximite avec Michelin en mix",
            "limitation": "Positionnement plus etroit, moins diversifie que Michelin",
            "market_cap": 6710,
            "enterprise_value": 7990,
            "ev_eqv": 7990 / 6710,
            "ev_sales": 1.18,
            "ev_ebitda": 4.96,
            "ev_ebit": 7.78,
            "ebit_margin": 0.1438,
        },
        {
            "company": "Yokohama Rubber",
            "ticker": "TYO:5101",
            "country": "Japon",
            "activity": "Pneus et caoutchouc industriel",
            "relevance": "Pneus et caoutchouc industriel, exposition Asie",
            "limitation": "Poids du caoutchouc industriel",
            "market_cap": 1070,
            "enterprise_value": 1500,
            "ev_eqv": 1500 / 1070,
            "ev_sales": 1.22,
            "ev_ebitda": 6.53,
            "ev_ebit": 9.60,
            "ebit_margin": 0.1267,
        },
        {
            "company": "Hankook Tire",
            "ticker": "KRX:161390",
            "country": "Coree du Sud",
            "activity": "Pneu mondial, bon mix valeur / premium",
            "relevance": "Pneus mondiaux, positionnement valeur/premium",
            "limitation": "Positionnement plus valeur, profil geographique different",
            "market_cap": 7430,
            "enterprise_value": 12940,
            "ev_eqv": 12940 / 7430,
            "ev_sales": 0.61,
            "ev_ebitda": 4.03,
            "ev_ebit": 7.07,
            "ebit_margin": 0.0868,
        },
    ]

    median_ev_sales = statistics.median(c["ev_sales"] for c in comps)
    median_ev_ebitda = statistics.median(c["ev_ebitda"] for c in comps)
    median_ev_ebit = statistics.median(c["ev_ebit"] for c in comps)
    median_ev_eqv = statistics.median(c["ev_eqv"] for c in comps)

    def percentile(values, p):
        values = sorted(values)
        pos = (len(values) - 1) * p
        lower = int(pos)
        upper = min(lower + 1, len(values) - 1)
        weight = pos - lower
        return values[lower] * (1 - weight) + values[upper] * weight

    comp_metrics = {
        "EV / CA": {
            "low": percentile([c["ev_sales"] for c in comps], 0.25),
            "mid": median_ev_sales,
            "high": percentile([c["ev_sales"] for c in comps], 0.75),
            "metric": forecast[2026]["sales"],
        },
        "EV / EBITDA": {
            "low": percentile([c["ev_ebitda"] for c in comps], 0.25),
            "mid": median_ev_ebitda,
            "high": percentile([c["ev_ebitda"] for c in comps], 0.75),
            "metric": forecast[2026]["ebitda"],
        },
        "EV / EBIT": {
            "low": percentile([c["ev_ebit"] for c in comps], 0.25),
            "mid": median_ev_ebit,
            "high": percentile([c["ev_ebit"] for c in comps], 0.75),
            "metric": forecast[2026]["ebit"],
        },
    }

    for item in comp_metrics.values():
        for key in ["low", "mid", "high"]:
            item[f"ev_{key}"] = item[key] * item["metric"]
            item[f"equity_{key}"] = item[f"ev_{key}"] - net_debt_2025
            item[f"price_{key}"] = item[f"equity_{key}"] / shares_m

    risk_free_rate = 0.0368
    erp = 0.0575
    beta = 1.00
    pre_tax_cost_debt = 0.0380
    cost_equity = risk_free_rate + beta * erp
    after_tax_cost_debt = pre_tax_cost_debt * (1 - tax_rate)
    market_cap = current_share_price * shares_m
    enterprise_value_market = market_cap + net_debt_2025
    weight_equity = market_cap / enterprise_value_market
    weight_debt = net_debt_2025 / enterprise_value_market
    wacc = weight_equity * cost_equity + weight_debt * after_tax_cost_debt
    terminal_growth = 0.015

    dcf = {}
    pv_fcff_total = 0
    for i, year in enumerate(forecast_years, start=1):
        discount_factor = 1 / ((1 + wacc) ** i)
        pv_fcff = forecast[year]["fcff"] * discount_factor
        dcf[year] = {
            **forecast[year],
            "discount_factor": discount_factor,
            "pv_fcff": pv_fcff,
        }
        pv_fcff_total += pv_fcff

    terminal_value = forecast[2030]["fcff"] * (1 + terminal_growth) / (wacc - terminal_growth)
    pv_terminal = terminal_value / ((1 + wacc) ** len(forecast_years))
    dcf_ev = pv_fcff_total + pv_terminal
    dcf_equity = dcf_ev - net_debt_2025
    dcf_price = dcf_equity / shares_m
    dcf_upside = dcf_price / current_share_price - 1

    sensitivity = {}
    wacc_range = [wacc - 0.010, wacc - 0.005, wacc, wacc + 0.005, wacc + 0.010]
    g_range = [0.005, 0.010, 0.015, 0.020, 0.025]
    for wr in wacc_range:
        sensitivity[wr] = {}
        for gr in g_range:
            pv_fcff = sum(forecast[y]["fcff"] / ((1 + wr) ** i) for i, y in enumerate(forecast_years, start=1))
            tv = forecast[2030]["fcff"] * (1 + gr) / (wr - gr)
            pv_tv = tv / ((1 + wr) ** len(forecast_years))
            price = (pv_fcff + pv_tv - net_debt_2025) / shares_m
            sensitivity[wr][gr] = price

    comps_mid_price = statistics.mean(item["price_mid"] for item in comp_metrics.values())
    blended_target = (comps_mid_price + dcf_price) / 2
    blended_upside = blended_target / current_share_price - 1
    total_return = blended_upside + 0.0426

    sales_2026_consensus = 26534
    segment_ebit_2026_consensus = 2982
    da_margin_2026_assumed = 0.073
    ebitda_2026_estimated = segment_ebit_2026_consensus + sales_2026_consensus * da_margin_2026_assumed

    return {
        "history": history,
        "forecast": forecast,
        "comps": comps,
        "comp_metrics": comp_metrics,
        "market": {
            "share_price": current_share_price,
            "shares_m": shares_m,
            "market_cap": market_cap,
            "enterprise_value_market": enterprise_value_market,
            "ev_eqv": enterprise_value_market / market_cap,
            "net_debt": net_debt_2025,
            "dividend_yield": 0.0426,
            "stockanalysis_ev": 24680,
        },
        "peer_stats": {
            "ev_eqv_low": percentile([c["ev_eqv"] for c in comps], 0.25),
            "median_ev_eqv": median_ev_eqv,
            "ev_eqv_high": percentile([c["ev_eqv"] for c in comps], 0.75),
        },
        "michelin_trading_multiples": {
            "ev_sales": enterprise_value_market / sales_2026_consensus,
            "ev_ebitda": enterprise_value_market / ebitda_2026_estimated,
            "ev_ebit": enterprise_value_market / segment_ebit_2026_consensus,
            "ev_eqv": enterprise_value_market / market_cap,
            "sales_2026_consensus": sales_2026_consensus,
            "segment_ebit_2026_consensus": segment_ebit_2026_consensus,
            "ebitda_2026_estimated": ebitda_2026_estimated,
            "da_margin_2026_assumed": da_margin_2026_assumed,
        },
        "wacc": {
            "risk_free_rate": risk_free_rate,
            "erp": erp,
            "beta": beta,
            "cost_equity": cost_equity,
            "pre_tax_cost_debt": pre_tax_cost_debt,
            "tax_rate": tax_rate,
            "after_tax_cost_debt": after_tax_cost_debt,
            "weight_equity": weight_equity,
            "weight_debt": weight_debt,
            "wacc": wacc,
            "terminal_growth": terminal_growth,
        },
        "dcf": {
            "years": forecast_years,
            "detail": dcf,
            "pv_fcff": pv_fcff_total,
            "terminal_value": terminal_value,
            "pv_terminal": pv_terminal,
            "enterprise_value": dcf_ev,
            "equity_value": dcf_equity,
            "price": dcf_price,
            "upside": dcf_upside,
            "sensitivity": sensitivity,
        },
        "recommendation": {
            "comps_mid_price": comps_mid_price,
            "blended_target": blended_target,
            "blended_upside": blended_upside,
            "total_return": total_return,
            "rating": "Acheter",
        },
    }


def style_sheet(ws):
    ws.sheet_view.showGridLines = False
    thin = Side(style="thin", color="D9DEE3")
    for row in ws.iter_rows():
        for cell in row:
            cell.alignment = Alignment(vertical="center")
            cell.border = Border(bottom=thin)
    for col in range(1, ws.max_column + 1):
        ws.column_dimensions[get_column_letter(col)].width = 16
    ws.column_dimensions["A"].width = 34


def header_cell(cell):
    cell.fill = PatternFill("solid", fgColor=BLUE)
    cell.font = Font(color="FFFFFF", bold=True)
    cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)


def section_cell(cell):
    cell.fill = PatternFill("solid", fgColor=LIGHT)
    cell.font = Font(color=DARK, bold=True)


def write_title(ws, title, subtitle=None):
    ws["A1"] = title
    ws["A1"].font = Font(size=16, bold=True, color=BLUE)
    if subtitle:
        ws["A2"] = subtitle
        ws["A2"].font = Font(size=10, italic=True, color=GRAY)


def write_excel(calcs: dict):
    wb = Workbook()
    wb.remove(wb.active)
    wb.calculation.fullCalcOnLoad = True

    # Sources
    ws = wb.create_sheet("Sources")
    write_title(ws, "Sources et conventions", "Projet finance Michelin, donnees collectees au 20 avril 2026")
    sources = [
        ("Michelin", "Financial information at December 31, 2025", "https://www.michelin.com/en/publications/group/financial-information-december-31-2025"),
        ("Michelin", "All Michelin key figures 2020-2025", "sources/michelin_key_figures_2020_2025.xlsx"),
        ("Michelin", "2025 Annual Results Guide", "sources/michelin_2025_annual_results_guide.pdf"),
        ("Michelin", "2025 Universal Registration Document", "sources/michelin_2025_urd.pdf"),
        ("Grand View Research", "Global tires market 2025-2030", "https://www.grandviewresearch.com/industry-analysis/tires-market-report"),
        ("StockAnalysis", "Michelin and peer valuation statistics", "https://stockanalysis.com/quote/epa/ML/statistics/"),
        ("CountryEconomy", "France 10-year bond yield", "https://countryeconomy.com/bonds/france"),
        ("Kroll", "Eurozone equity risk premium", "https://www.kroll.com/en/reports/cost-of-capital/recommended-eurozone-equity-risk-premium-corresponding-risk-free-rates"),
    ]
    ws.append([])
    ws.append(["Source", "Donnee utilisee", "Lien / fichier"])
    for cell in ws[3]:
        header_cell(cell)
    for row in sources:
        ws.append(row)
    ws.append([])
    ws.append(["Convention", "Valeur"])
    for cell in ws[ws.max_row]:
        header_cell(cell)
    ws.append(["Unite des etats financiers", "EUR millions sauf indication contraire"])
    ws.append(["Cours de reference Michelin", calcs["market"]["share_price"]])
    ws.append(["Actions diluees / en circulation", f"{calcs['market']['shares_m']:.2f}m"])
    ws.append(["Recommandation retenue", calcs["recommendation"]["rating"]])
    ws.freeze_panes = "A4"
    style_sheet(ws)

    # Dashboard
    ws = wb.create_sheet("Dashboard")
    write_title(ws, "Dashboard Michelin", "Lecture rapide du modele: valorisation, message cle et feuilles a utiliser")
    ws.append([])
    ws.append(["Indicateur", "Valeur"])
    for cell in ws[3]:
        header_cell(cell)
    dashboard_rows = [
        ("Cours actuel", "=WACC!B5"),
        ("Valeur comparables", "=AVERAGE(Comps!F13:F15)"),
        ("Valeur DCF", "=DCF!B23"),
        ("Objectif central", "=Football_Field!C8"),
        ("Rendement dividende", "=0.0426"),
        ("Rendement total", "=((B7/B4)-1)+B8"),
    ]
    for r, (label, value) in enumerate(dashboard_rows, start=4):
        ws.cell(r, 1, label)
        ws.cell(r, 2, value)
        section_cell(ws.cell(r, 1))
    for r in [4, 5, 6, 7]:
        ws.cell(r, 2).number_format = 'EUR 0.0'
    for r in [8, 9]:
        ws.cell(r, 2).number_format = "0.0%"

    ws["D4"] = "Comment lire le fichier"
    header_cell(ws["D4"])
    guidance = [
        "1. Modifier les hypotheses de croissance et de marge dans Forecast.",
        "2. Modifier les inputs de marche et de WACC dans WACC.",
        "3. Les feuilles DCF, Sensitivity, Tornado et Football_Field se mettent ensuite a jour.",
        "4. Scenario_DCF et Buyback_Impact servent a tester des lectures alternatives sans casser la base.",
    ]
    for i, text in enumerate(guidance, start=5):
        ws.cell(i, 4, text)
        ws.cell(i, 4).alignment = Alignment(wrap_text=True, vertical="top")

    ws["G4"] = "Lecture"
    ws["H4"] = "Valeur"
    header_cell(ws["G4"])
    header_cell(ws["H4"])
    chart_labels = ["Cours actuel", "Comparables", "DCF", "Objectif"]
    chart_formulas = ["=B4", "=B5", "=B6", "=B7"]
    for i, (label, formula) in enumerate(zip(chart_labels, chart_formulas), start=5):
        ws.cell(i, 7, label)
        ws.cell(i, 8, formula)
        ws.cell(i, 8).number_format = 'EUR 0.0'
    chart = BarChart()
    chart.title = "Synthese de valorisation"
    chart.y_axis.title = "EUR/action"
    data = Reference(ws, min_col=8, max_col=8, min_row=4, max_row=8)
    cats = Reference(ws, min_col=7, max_col=7, min_row=5, max_row=8)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    chart.height = 6
    chart.width = 8
    ws.add_chart(chart, "J4")
    style_sheet(ws)

    # Historical
    ws = wb.create_sheet("Historical")
    years = calcs["history"]["years"]
    hist = calcs["history"]["rows"]
    write_title(ws, "Historique financier Michelin", "Source: fichier officiel Michelin key figures 2020-2025")
    ws.append([])
    ws.append(["EURm"] + years)
    for cell in ws[3]:
        header_cell(cell)
    rows = [
        ("Chiffre d'affaires", "sales", "number"),
        ("Croissance CA", "sales_growth", "percent"),
        ("Segment EBITDA", "ebitda", "number"),
        ("Marge EBITDA", None, "percent_formula"),
        ("Resultat operationnel des secteurs", "segment_ebit", "number"),
        ("Marge operationnelle secteurs", "segment_margin", "percent"),
        ("Resultat operationnel", "operating_income", "number"),
        ("Resultat net", "net_income", "number"),
        ("Marge nette", None, "percent_formula"),
        ("Cash-flow operationnel", "cfo", "number"),
        ("Capex", "capex", "number"),
        ("Capex / CA", None, "percent_formula_capex"),
        ("Free cash-flow", "fcf", "number"),
        ("FCF / CA", None, "percent_formula_fcf"),
        ("Dette nette", "net_debt", "number"),
        ("Dette nette / EBITDA", None, "number_formula_leverage"),
    ]
    start_row = 4
    for r, (label, key, row_type) in enumerate(rows, start=start_row):
        ws.cell(r, 1, label)
        for c, year in enumerate(years, start=2):
            if key:
                ws.cell(r, c, hist[key][year])
            elif row_type == "percent_formula":
                ws.cell(r, c, f"={get_column_letter(c)}6/{get_column_letter(c)}4")
            elif row_type == "percent_formula_capex":
                ws.cell(r, c, f"={get_column_letter(c)}14/{get_column_letter(c)}4")
            elif row_type == "percent_formula_fcf":
                ws.cell(r, c, f"={get_column_letter(c)}16/{get_column_letter(c)}4")
            elif row_type == "number_formula_leverage":
                ws.cell(r, c, f"={get_column_letter(c)}18/{get_column_letter(c)}6")
        section_cell(ws.cell(r, 1))
    for row in ws.iter_rows(min_row=start_row, max_row=start_row + len(rows) - 1, min_col=2, max_col=7):
        for cell in row:
            if cell.row in [5, 7, 12, 15, 17]:
                cell.number_format = "0.0%"
            elif cell.row == 19:
                cell.number_format = "0.0x"
            else:
                cell.number_format = '#,##0'
    chart = BarChart()
    chart.title = "Chiffre d'affaires et EBIT secteurs"
    chart.y_axis.title = "EURm"
    data = Reference(ws, min_col=2, max_col=7, min_row=4, max_row=8)
    cats = Reference(ws, min_col=2, max_col=7, min_row=3)
    chart.add_data(data, titles_from_data=False)
    chart.set_categories(cats)
    chart.height = 7
    chart.width = 15
    ws.add_chart(chart, "I4")
    style_sheet(ws)

    # Forecast
    ws = wb.create_sheet("Forecast")
    write_title(ws, "Projection 2026E-2030E", "Scenario central: reprise graduelle des volumes, mix premium et discipline cash")
    forecast_years = list(calcs["forecast"].keys())
    ws.append([])
    ws.append(["EURm"] + [f"{y}E" for y in forecast_years])
    for cell in ws[3]:
        header_cell(cell)
    forecast_labels = [
        "Chiffre d'affaires",
        "Croissance CA",
        "Marge EBITDA",
        "EBITDA",
        "Marge EBIT secteurs",
        "EBIT secteurs",
        "Taux d'impot",
        "NOPAT",
        "D&A / CA",
        "D&A",
        "Capex / CA",
        "Capex",
        "Variation BFR",
        "FCFF",
    ]
    for r, label in enumerate(forecast_labels, start=4):
        ws.cell(r, 1, label)
        section_cell(ws.cell(r, 1))
    for c, year in enumerate(forecast_years, start=2):
        col = get_column_letter(c)
        prev_col = get_column_letter(c - 1)
        if c == 2:
            ws.cell(4, c, f"=Historical!G4*(1+{col}5)")
            ws.cell(16, c, f"=({col}4-Historical!G4)*0.01")
        else:
            ws.cell(4, c, f"={prev_col}4*(1+{col}5)")
            ws.cell(16, c, f"=({col}4-{prev_col}4)*0.01")
        ws.cell(5, c, calcs["forecast"][year]["growth"])
        ws.cell(6, c, calcs["forecast"][year]["ebitda_margin"])
        ws.cell(7, c, f"={col}4*{col}6")
        ws.cell(8, c, calcs["forecast"][year]["ebit_margin"])
        ws.cell(9, c, f"={col}4*{col}8")
        ws.cell(10, c, calcs["forecast"][year]["tax_rate"])
        ws.cell(11, c, f"={col}9*(1-{col}10)")
        ws.cell(12, c, calcs["forecast"][year]["da_margin"])
        ws.cell(13, c, f"={col}4*{col}12")
        ws.cell(14, c, calcs["forecast"][year]["capex_margin"])
        ws.cell(15, c, f"={col}4*{col}14")
        ws.cell(17, c, f"={col}11+{col}13-{col}15-{col}16")
    for row_idx in [5, 6, 8, 10, 12, 14]:
        for c in range(2, 7):
            ws.cell(row_idx, c).number_format = "0.0%"
    for row_idx in [4, 7, 9, 11, 13, 15, 16, 17]:
        for c in range(2, 7):
            ws.cell(row_idx, c).number_format = '#,##0'
    chart = LineChart()
    chart.title = "Projection du CA et du FCFF"
    chart.y_axis.title = "EURm"
    data = Reference(ws, min_col=2, max_col=6, min_row=4, max_row=17)
    cats = Reference(ws, min_col=2, max_col=6, min_row=3)
    chart.add_data(data, titles_from_data=False)
    chart.set_categories(cats)
    chart.height = 7
    chart.width = 15
    ws.add_chart(chart, "I4")
    style_sheet(ws)

    # Comps
    ws = wb.create_sheet("Comps")
    write_title(ws, "Comparables boursiers", "Multiples de marche ponctuels au 20 avril 2026 ou derniere cloture disponible")
    ws.append([])
    ws.append(["Societe", "Ticker", "Pays", "Justification", "EV/CA", "EV/EBITDA", "EV/EBIT"])
    for cell in ws[3]:
        header_cell(cell)
    for comp in calcs["comps"]:
        ws.append([
            comp["company"],
            comp["ticker"],
            comp["country"],
            comp["relevance"],
            comp["ev_sales"],
            comp["ev_ebitda"],
            comp["ev_ebit"],
        ])
    for r in range(4, 10):
        for c in range(5, 8):
            ws.cell(r, c).number_format = "0.00x"
    ws.append([])
    start = ws.max_row + 1
    ws.append(["Statistique", "EV/CA", "EV/EBITDA", "EV/EBIT"])
    for cell in ws[start]:
        header_cell(cell)
    ws.append(["Quartile bas", "=QUARTILE.INC(E4:E9,1)", "=QUARTILE.INC(F4:F9,1)", "=QUARTILE.INC(G4:G9,1)"])
    ws.append(["Mediane", "=MEDIAN(E4:E9)", "=MEDIAN(F4:F9)", "=MEDIAN(G4:G9)"])
    ws.append(["Quartile haut", "=QUARTILE.INC(E4:E9,3)", "=QUARTILE.INC(F4:F9,3)", "=QUARTILE.INC(G4:G9,3)"])
    for r in range(start + 1, start + 4):
        for c in range(2, 5):
            ws.cell(r, c).number_format = "0.00x"
    ws.append([])
    val_start = ws.max_row + 1
    ws.append(["Application 2026E", "Multiple bas", "Multiple median", "Multiple haut", "Prix bas", "Prix median", "Prix haut"])
    for cell in ws[val_start]:
        header_cell(cell)
    for method, item in calcs["comp_metrics"].items():
        ws.append([
            method,
            item["low"],
            item["mid"],
            item["high"],
            item["price_low"],
            item["price_mid"],
            item["price_high"],
        ])
    for r in range(val_start + 1, val_start + 4):
        for c in range(2, 5):
            ws.cell(r, c).number_format = "0.00x"
        for c in range(5, 8):
            ws.cell(r, c).number_format = 'EUR 0.0'
    style_sheet(ws)
    ws.column_dimensions["D"].width = 46

    # WACC
    ws = wb.create_sheet("WACC")
    write_title(ws, "Calcul du WACC", "Approche CAPM et structure de capital de marche")
    inputs = [
        ("Cours Michelin", calcs["market"]["share_price"], "EUR/action"),
        ("Actions en circulation", calcs["market"]["shares_m"], "m"),
        ("Capitalisation boursiere", calcs["market"]["market_cap"], "EURm"),
        ("Dette nette", calcs["market"]["net_debt"], "EURm"),
        ("Enterprise value de marche", calcs["market"]["enterprise_value_market"], "EURm"),
        ("Taux sans risque France 10 ans", calcs["wacc"]["risk_free_rate"], "%"),
        ("Prime de risque actions Eurozone", calcs["wacc"]["erp"], "%"),
        ("Beta Michelin", calcs["wacc"]["beta"], "x"),
        ("Cout des fonds propres", calcs["wacc"]["cost_equity"], "%"),
        ("Cout de la dette avant impots", calcs["wacc"]["pre_tax_cost_debt"], "%"),
        ("Taux d'impot effectif", calcs["wacc"]["tax_rate"], "%"),
        ("Cout de la dette apres impots", calcs["wacc"]["after_tax_cost_debt"], "%"),
        ("Poids fonds propres", calcs["wacc"]["weight_equity"], "%"),
        ("Poids dette nette", calcs["wacc"]["weight_debt"], "%"),
        ("WACC retenu", calcs["wacc"]["wacc"], "%"),
        ("Croissance terminale", calcs["wacc"]["terminal_growth"], "%"),
    ]
    ws.append([])
    ws.append(["Input", "Valeur", "Unite"])
    for cell in ws[3]:
        header_cell(cell)
    for row in inputs:
        ws.append(row)
    for r in range(4, 20):
        unit = ws.cell(r, 3).value
        if unit == "%":
            ws.cell(r, 2).number_format = "0.0%"
        elif unit == "x":
            ws.cell(r, 2).number_format = "0.00x"
        else:
            ws.cell(r, 2).number_format = '#,##0.0'
    style_sheet(ws)

    # DCF
    ws = wb.create_sheet("DCF")
    write_title(ws, "Valorisation DCF", "FCFF 2026E-2030E + valeur terminale par croissance perpetuelle")
    years = calcs["dcf"]["years"]
    ws.append([])
    ws.append(["EURm"] + years)
    for cell in ws[3]:
        header_cell(cell)
    dcf_labels = [
        "Chiffre d'affaires",
        "EBIT secteurs",
        "NOPAT",
        "D&A",
        "Capex",
        "Variation BFR",
        "FCFF",
        "Facteur d'actualisation",
        "PV FCFF",
    ]
    for r, label in enumerate(dcf_labels, start=4):
        ws.cell(r, 1, label)
        section_cell(ws.cell(r, 1))
    forecast_row_map = {4: 4, 5: 9, 6: 11, 7: 13, 8: 15, 9: 16, 10: 17}
    for c, _year in enumerate(years, start=2):
        col = get_column_letter(c)
        for dcf_row, forecast_row in forecast_row_map.items():
            ws.cell(dcf_row, c, f"=Forecast!{col}{forecast_row}")
        ws.cell(11, c, "=(1+WACC!$B$19)^-(COLUMN()-1)")
        ws.cell(12, c, f"={col}10*{col}11")
    for r in [11]:
        for c in range(2, 7):
            ws.cell(r, c).number_format = "0.000x"
    for r in [4, 5, 6, 7, 8, 9, 10, 12]:
        for c in range(2, 7):
            ws.cell(r, c).number_format = '#,##0'
    summary_start = 15
    summary = [
        ("PV des FCFF explicites", "=SUM(B12:F12)"),
        ("Valeur terminale", "=F10*(1+WACC!$B$20)/(WACC!$B$19-WACC!$B$20)"),
        ("PV valeur terminale", "=B17*F11"),
        ("Enterprise value DCF", "=B16+B18"),
        ("Dette nette", "=WACC!B8"),
        ("Equity value DCF", "=B19-B20"),
        ("Actions en circulation", "=WACC!B6"),
        ("Valeur par action DCF", "=B21/B22"),
        ("Cours actuel", "=WACC!B5"),
        ("Potentiel DCF", "=B23/B24-1"),
    ]
    ws.cell(summary_start, 1, "Synthese")
    header_cell(ws.cell(summary_start, 1))
    for i, (label, value) in enumerate(summary, start=summary_start + 1):
        ws.cell(i, 1, label)
        ws.cell(i, 2, value)
        if "Potentiel" in label:
            ws.cell(i, 2).number_format = "0.0%"
        elif "Valeur par action" in label or "Cours" in label:
            ws.cell(i, 2).number_format = 'EUR 0.0'
        else:
            ws.cell(i, 2).number_format = '#,##0.0'
    style_sheet(ws)

    # Scenario DCF
    ws = wb.create_sheet("Scenario_DCF")
    write_title(ws, "Scenario DCF", "Trois lectures du DCF avec facade simple et calculs detailles caches")
    ws.append([])
    ws.append(["Hypothese", "Baissier", "Central", "Haussier"])
    for cell in ws[3]:
        header_cell(cell)
    scenario_inputs = [
        ("Croissance 2027E", 0.000, "=Forecast!C5", 0.025),
        ("Croissance 2028E", 0.005, "=Forecast!D5", 0.030),
        ("Croissance 2029E", 0.010, "=Forecast!E5", 0.030),
        ("Croissance 2030E", 0.015, "=Forecast!F5", 0.030),
        ("Marge EBIT 2030", 0.105, "=Forecast!F8", 0.130),
        ("WACC", "=WACC!B19+0.005", "=WACC!B19", "=WACC!B19-0.005"),
        ("Croissance terminale", "=WACC!B20-0.005", "=WACC!B20", "=WACC!B20+0.005"),
    ]
    for r, (label, bear, base, bull) in enumerate(scenario_inputs, start=4):
        ws.cell(r, 1, label)
        ws.cell(r, 2, bear)
        ws.cell(r, 3, base)
        ws.cell(r, 4, bull)
        section_cell(ws.cell(r, 1))
        for c in range(2, 5):
            ws.cell(r, c).number_format = "0.0%"
    ws["A13"] = "Sortie"
    ws["B13"] = "Baissier"
    ws["C13"] = "Central"
    ws["D13"] = "Haussier"
    for cell in ws[13]:
        if cell.column <= 4:
            header_cell(cell)
    scenario_cols = {2: 8, 3: 15, 4: 22}
    output_rows = {
        "EV DCF": 19,
        "EqV DCF": 20,
        "Valeur / action": 21,
        "Potentiel vs cours": 22,
    }
    for row_idx, label in enumerate(output_rows.keys(), start=14):
        ws.cell(row_idx, 1, label)
        section_cell(ws.cell(row_idx, 1))
    for input_col, start_col in scenario_cols.items():
        start_letter = get_column_letter(start_col)
        year_letters = [get_column_letter(start_col + i) for i in range(5)]
        ws.cell(4, start_col, {2: "Baissier", 3: "Central", 4: "Haussier"}[input_col])
        for i, year in enumerate([2026, 2027, 2028, 2029, 2030], start=start_col):
            ws.cell(5, i, f"{year}E")
        ws.cell(6, start_col, "=Forecast!B4")
        ws.cell(7, start_col, "=Forecast!B5")
        ws.cell(8, start_col, "=Forecast!B8")
        for offset, year_col in enumerate(range(start_col + 1, start_col + 5), start=1):
            col = get_column_letter(year_col)
            prev = get_column_letter(year_col - 1)
            growth_row = 4 + offset
            ws.cell(6, year_col, f"={prev}6*(1+${get_column_letter(input_col)}${growth_row})")
            ws.cell(7, year_col, f"=${get_column_letter(input_col)}${growth_row}")
            ws.cell(8, year_col, f"=Forecast!$B$8+((${get_column_letter(input_col)}$9-Forecast!$B$8)*{offset}/4)")
        for col in year_letters:
            ws[f"{col}9"] = f"={col}6*{col}8"
            ws[f"{col}10"] = f"={col}9*(1-WACC!$B$15)"
            forecast_col = get_column_letter(2 + (ord(col) - ord(start_letter)))
            ws[f"{col}11"] = f"={col}6*Forecast!{forecast_col}12"
            ws[f"{col}12"] = f"={col}6*Forecast!{forecast_col}14"
        ws[f"{start_letter}13"] = f"=({start_letter}6-Historical!G4)*0.01"
        for idx, year_col in enumerate(year_letters[1:], start=1):
            prev = year_letters[idx - 1]
            ws[f"{year_col}13"] = f"=({year_col}6-{prev}6)*0.01"
        for col in year_letters:
            ws[f"{col}14"] = f"={col}10+{col}11-{col}12-{col}13"
            ws[f"{col}15"] = f"=(1-${get_column_letter(input_col)}$10)^0"
        for idx, col in enumerate(year_letters, start=1):
            ws[f"{col}15"] = f"=1/((1+${get_column_letter(input_col)}$10)^{idx})"
            ws[f"{col}16"] = f"={col}14*{col}15"
        last_col = year_letters[-1]
        ws[f"{start_letter}17"] = f"={last_col}14*(1+${get_column_letter(input_col)}$11)/(${get_column_letter(input_col)}$10-${get_column_letter(input_col)}$11)"
        ws[f"{start_letter}18"] = f"={start_letter}17*{last_col}15"
        ws[f"{start_letter}19"] = f"=SUM({start_letter}16:{last_col}16)+{start_letter}18"
        ws[f"{start_letter}20"] = f"={start_letter}19-WACC!$B$8"
        ws[f"{start_letter}21"] = f"={start_letter}20/WACC!$B$6"
        ws[f"{start_letter}22"] = f"={start_letter}21/WACC!$B$5-1"
        visible_col = get_column_letter(input_col)
        ws[f"{visible_col}14"] = f"={start_letter}19"
        ws[f"{visible_col}15"] = f"={start_letter}20"
        ws[f"{visible_col}16"] = f"={start_letter}21"
        ws[f"{visible_col}17"] = f"={start_letter}22"
    for c in range(2, 5):
        for r in [14, 15, 16]:
            ws.cell(r, c).number_format = 'EUR 0.0'
        ws.cell(17, c).number_format = "0.0%"
    ws["F4"] = "Lecture"
    header_cell(ws["F4"])
    ws["F5"] = "2026E reste ancre sur la base Forecast."
    ws["F6"] = "Le scenario agit sur la croissance 2027E-2030E, la marge 2030, le WACC et g."
    ws["F7"] = "Les calculs detailles sont caches a droite pour garder une feuille lisible."
    chart = BarChart()
    chart.title = "Valeur / action par scenario"
    chart.y_axis.title = "EUR/action"
    data = Reference(ws, min_col=2, max_col=4, min_row=16, max_row=16)
    cats = Reference(ws, min_col=2, max_col=4, min_row=13, max_row=13)
    chart.add_data(data, titles_from_data=False, from_rows=True)
    chart.set_categories(cats)
    chart.height = 6
    chart.width = 8
    ws.add_chart(chart, "F9")
    for col_idx in range(8, 28):
        ws.column_dimensions[get_column_letter(col_idx)].hidden = True
    style_sheet(ws)

    # Sensitivity
    ws = wb.create_sheet("Sensitivity")
    write_title(ws, "Sensibilite DCF", "Valeur par action selon WACC et croissance terminale")
    ws.append([])
    ws.append(["WACC \\ g", "g -100 bps", "g -50 bps", "g central", "g +50 bps", "g +100 bps"])
    for cell in ws[3]:
        header_cell(cell)
    g_offsets = [-0.010, -0.005, 0.0, 0.005, 0.010]
    wacc_offsets = [-0.010, -0.005, 0.0, 0.005, 0.010]
    for c, offset in enumerate(g_offsets, start=2):
        col = get_column_letter(c)
        if offset == 0:
            ws.cell(4, c, "=WACC!$B$20")
        elif offset > 0:
            ws.cell(4, c, f"=WACC!$B$20+{offset}")
        else:
            ws.cell(4, c, f"=WACC!$B$20{offset}")
        ws.cell(4, c).number_format = "0.0%"
    for r, offset in enumerate(wacc_offsets, start=5):
        if offset == 0:
            ws.cell(r, 1, "=WACC!$B$19")
        elif offset > 0:
            ws.cell(r, 1, f"=WACC!$B$19+{offset}")
        else:
            ws.cell(r, 1, f"=WACC!$B$19{offset}")
        ws.cell(r, 1).number_format = "0.0%"
        for c in range(2, 7):
            col = get_column_letter(c)
            ws.cell(
                r,
                c,
                (
                    f"=((DCF!$B$10/(1+$A{r})^1)"
                    f"+(DCF!$C$10/(1+$A{r})^2)"
                    f"+(DCF!$D$10/(1+$A{r})^3)"
                    f"+(DCF!$E$10/(1+$A{r})^4)"
                    f"+(DCF!$F$10/(1+$A{r})^5)"
                    f"+((DCF!$F$10*(1+{col}$4))/($A{r}-{col}$4))/((1+$A{r})^5)"
                    f"-WACC!$B$8)/WACC!$B$6"
                ),
            )
            ws.cell(r, c).number_format = 'EUR 0.0'
    style_sheet(ws)

    # Tornado
    ws = wb.create_sheet("Tornado")
    write_title(ws, "Tornado", "Impact des hypotheses cles sur la valeur DCF par action")
    ws.append([])
    ws.append(["Variable", "Hypothese basse", "Hypothese haute", "Valeur basse", "Valeur haute", "Amplitude"])
    for cell in ws[3]:
        header_cell(cell)
    tornado_rows = [
        ("Marge EBIT", "-100 bps", "+100 bps"),
        ("WACC", "+50 bps", "-50 bps"),
        ("Capex / CA", "+50 bps", "-50 bps"),
        ("Croissance terminale", "-50 bps", "+50 bps"),
        ("Croissance CA", "-100 bps", "+100 bps"),
    ]
    for r, (label, low_case, high_case) in enumerate(tornado_rows, start=4):
        ws.cell(r, 1, label)
        ws.cell(r, 2, low_case)
        ws.cell(r, 3, high_case)
        section_cell(ws.cell(r, 1))
    helper_specs = [
        ("sales_low", "sales", -0.010),
        ("sales_high", "sales", 0.010),
        ("margin_low", "margin", -0.010),
        ("margin_high", "margin", 0.010),
        ("capex_low", "capex", -0.005),
        ("capex_high", "capex", 0.005),
    ]
    helper_price_refs = {}
    for idx, (name, driver, offset) in enumerate(helper_specs):
        start_col = 8 + idx * 6
        letters = [get_column_letter(start_col + i) for i in range(5)]
        helper_price_refs[name] = f"{letters[0]}30"
        ws.cell(20, start_col, name)
        for i, year in enumerate([2026, 2027, 2028, 2029, 2030], start=start_col):
            ws.cell(21, i, f"{year}E")
        for j, col in enumerate(letters):
            forecast_col = get_column_letter(2 + j)
            if driver == "sales":
                if j == 0:
                    ws[f"{col}22"] = f"=Historical!G4*(1+(Forecast!{forecast_col}5{offset:+.3f}))"
                else:
                    prev = letters[j - 1]
                    ws[f"{col}22"] = f"={prev}22*(1+(Forecast!{forecast_col}5{offset:+.3f}))"
            else:
                ws[f"{col}22"] = f"=Forecast!{forecast_col}4"
            if driver == "margin":
                ws[f"{col}23"] = f"=Forecast!{forecast_col}8{offset:+.3f}"
            else:
                ws[f"{col}23"] = f"=Forecast!{forecast_col}8"
            ws[f"{col}24"] = f"={col}22*{col}23"
            ws[f"{col}25"] = f"={col}24*(1-WACC!$B$15)"
            ws[f"{col}26"] = f"={col}22*Forecast!{forecast_col}12"
            if driver == "capex":
                ws[f"{col}27"] = f"={col}22*(Forecast!{forecast_col}14{offset:+.3f})"
            else:
                ws[f"{col}27"] = f"={col}22*Forecast!{forecast_col}14"
            if j == 0:
                ws[f"{col}28"] = f"=({col}22-Historical!G4)*0.01"
            else:
                prev = letters[j - 1]
                ws[f"{col}28"] = f"=({col}22-{prev}22)*0.01"
            ws[f"{col}29"] = f"={col}25+{col}26-{col}27-{col}28"
        last_col = letters[-1]
        first_col = letters[0]
        ws[f"{first_col}30"] = (
            f"=(({letters[0]}29*DCF!B11)+({letters[1]}29*DCF!C11)+({letters[2]}29*DCF!D11)"
            f"+({letters[3]}29*DCF!E11)+({letters[4]}29*DCF!F11)"
            f"+(({last_col}29*(1+WACC!$B$20)/(WACC!$B$19-WACC!$B$20))*DCF!F11)-WACC!$B$8)/WACC!$B$6"
        )
    ws["D4"] = f"={helper_price_refs['margin_low']}"
    ws["E4"] = f"={helper_price_refs['margin_high']}"
    ws["D5"] = "=Sensitivity!D8"
    ws["E5"] = "=Sensitivity!D6"
    ws["D6"] = f"={helper_price_refs['capex_high']}"
    ws["E6"] = f"={helper_price_refs['capex_low']}"
    ws["D7"] = "=Sensitivity!C7"
    ws["E7"] = "=Sensitivity!E7"
    ws["D8"] = f"={helper_price_refs['sales_low']}"
    ws["E8"] = f"={helper_price_refs['sales_high']}"
    for r in range(4, 9):
        ws.cell(r, 6, f"=E{r}-D{r}")
        for c in range(4, 7):
            ws.cell(r, c).number_format = 'EUR 0.0'
    chart = BarChart()
    chart.type = "bar"
    chart.title = "Sensibilite des variables cles"
    chart.x_axis.title = "EUR/action"
    data = Reference(ws, min_col=4, max_col=5, min_row=3, max_row=8)
    cats = Reference(ws, min_col=1, max_col=1, min_row=4, max_row=8)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    chart.height = 6
    chart.width = 8
    ws.add_chart(chart, "H4")
    for col_idx in range(8, 44):
        ws.column_dimensions[get_column_letter(col_idx)].hidden = True
    style_sheet(ws)

    # Buyback impact
    ws = wb.create_sheet("Buyback_Impact")
    write_title(ws, "Buyback Impact", "Accretion par action si Michelin rachete ses titres sous la valeur intrinseque")
    ws.append([])
    ws.append(["Input", "Valeur"])
    for cell in ws[3]:
        header_cell(cell)
    buyback_inputs = [
        ("Equity value DCF de base", "=DCF!B21"),
        ("Actions de base", "=WACC!B6"),
        ("Valeur intrinseque / action", "=DCF!B23"),
        ("Prix moyen de rachat", "=WACC!B5"),
        ("Programme minimum 2026-2028", 2000),
    ]
    for r, (label, value) in enumerate(buyback_inputs, start=4):
        ws.cell(r, 1, label)
        ws.cell(r, 2, value)
        section_cell(ws.cell(r, 1))
    for r in [4, 6, 7]:
        ws.cell(r, 2).number_format = 'EUR 0.0'
    ws.cell(5, 2).number_format = '#,##0.0'
    ws.cell(8, 2).number_format = '#,##0'

    ws["A11"] = "Programme (EURm)"
    ws["B11"] = "Actions retirees"
    ws["C11"] = "Actions post-rachat"
    ws["D11"] = "EqV post-rachat"
    ws["E11"] = "Valeur/action"
    ws["F11"] = "Accretion"
    for c in range(1, 7):
        header_cell(ws.cell(11, c))
    amounts = [0, 500, 1000, 1500, 2000, 3000, 4000]
    for i, amount in enumerate(amounts, start=12):
        ws.cell(i, 1, amount)
        ws.cell(i, 2, f"=A{i}/$B$7")
        ws.cell(i, 3, f"=$B$5-B{i}")
        ws.cell(i, 4, f"=$B$4-A{i}")
        ws.cell(i, 5, f"=D{i}/C{i}")
        ws.cell(i, 6, f"=E{i}/$B$6-1")
        ws.cell(i, 1).number_format = '#,##0'
        ws.cell(i, 2).number_format = '#,##0.0'
        ws.cell(i, 3).number_format = '#,##0.0'
        ws.cell(i, 4).number_format = 'EUR 0.0'
        ws.cell(i, 5).number_format = 'EUR 0.0'
        ws.cell(i, 6).number_format = '0.0%'
    ws["H4"] = "Lecture"
    header_cell(ws["H4"])
    ws["H5"] = "Le buyback detruit d'abord de l'equity value du montant depense."
    ws["H6"] = "Il cree de la valeur par action seulement si le prix de rachat est inferieur a la valeur intrinseque."
    ws["H7"] = "Le tableau montre l'accretion theoriquement impliquee pour plusieurs tailles de programme."
    chart = LineChart()
    chart.title = "Valeur par action apres rachat"
    chart.y_axis.title = "EUR/action"
    data = Reference(ws, min_col=5, max_col=5, min_row=11, max_row=18)
    cats = Reference(ws, min_col=1, max_col=1, min_row=12, max_row=18)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    chart.height = 6
    chart.width = 8
    ws.add_chart(chart, "H10")
    style_sheet(ws)

    # Football field
    ws = wb.create_sheet("Football_Field")
    write_title(ws, "Football field", "Synthese des valorisations par action")
    ws.append([])
    ws.append(["Methode", "Bas", "Central", "Haut", "Commentaire"])
    for cell in ws[3]:
        header_cell(cell)
    ws.append(["Cours actuel", "=WACC!B5", "=WACC!B5", "=WACC!B5", "Reference marche StockAnalysis, 20 avril 2026"])
    ws.append(["Comparables", "=MIN(Comps!E13:E15)", "=AVERAGE(Comps!F13:F15)", "=MAX(Comps!G13:G15)", "Application quartiles / mediane aux metrics 2026E"])
    ws.append(["DCF", "=MIN(Sensitivity!B5:F9)", "=DCF!B23", "=MAX(Sensitivity!B5:F9)", "Sensibilite WACC et croissance terminale"])
    ws.append(["Objectif retenu", "=AVERAGE(C6,C7)", "=AVERAGE(C6,C7)", "=AVERAGE(C6,C7)", "Moyenne du DCF central et des comparables centraux"])
    for r in range(4, 8):
        for c in range(2, 5):
            ws.cell(r, c).number_format = 'EUR 0.0'
    style_sheet(ws)
    ws.column_dimensions["E"].width = 48

    for sheet in wb.worksheets:
        sheet.freeze_panes = "A4"
        for row in sheet.iter_rows():
            for cell in row:
                cell.alignment = Alignment(vertical="center", wrap_text=True)

    wb.calculation.calcMode = "auto"
    wb.calculation.fullCalcOnLoad = True
    wb.calculation.forceFullCalc = True
    wb.save(EXCEL_OUT)


def add_slide_title(slide, title, subtitle=None):
    shape = slide.shapes.add_textbox(Inches(0.35), Inches(0.22), Inches(12.6), Inches(0.55))
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(BLUE)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.38), Inches(0.78), Inches(12.3), Inches(0.3))
        tf2 = sub.text_frame
        tf2.text = subtitle
        tf2.paragraphs[0].font.size = Pt(9)
        tf2.paragraphs[0].font.color.rgb = RGBColor.from_string(GRAY)


def add_footer(slide, text):
    shape = slide.shapes.add_textbox(Inches(0.35), Inches(7.05), Inches(12.6), Inches(0.25))
    tf = shape.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(7.5)
    p.font.color.rgb = RGBColor.from_string(GRAY)


def add_bullets(slide, x, y, w, h, bullets, font_size=13):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor.from_string(DARK)
        p.space_after = Pt(5)
        p.line_spacing = 1.05
    return shape


def add_metric(slide, x, y, w, h, value, label, color=TEAL):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor.from_string(LIGHT)
    box.line.color.rgb = RGBColor.from_string(MID)
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(color)
    p2 = tf.add_paragraph()
    p2.text = label
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(9)
    p2.font.color.rgb = RGBColor.from_string(DARK)


def add_table(slide, x, y, w, h, rows, col_widths=None, font_size=8.5):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(font_size)
            para.font.color.rgb = RGBColor.from_string(DARK)
            para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(BLUE)
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string("FFFFFF" if r % 2 else LIGHT)
    return table_shape


def add_horizontal_bars(slide, x, y, w, h, data, value_max, colors=None):
    colors = colors or [TEAL, BLUE, GOLD, RED]
    row_h = h / len(data)
    for i, (label, value) in enumerate(data):
        yy = y + i * row_h
        lbl = slide.shapes.add_textbox(Inches(x), Inches(yy + 0.05), Inches(2.2), Inches(row_h - 0.08))
        lbl.text_frame.text = label
        lbl.text_frame.paragraphs[0].font.size = Pt(9)
        lbl.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(DARK)
        bar_w = max(0.05, (value / value_max) * (w - 3.0))
        bar = slide.shapes.add_shape(1, Inches(x + 2.3), Inches(yy + 0.09), Inches(bar_w), Inches(row_h - 0.18))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(colors[i % len(colors)])
        bar.line.fill.background()
        val = slide.shapes.add_textbox(Inches(x + 2.35 + bar_w), Inches(yy + 0.05), Inches(0.9), Inches(row_h - 0.08))
        val.text_frame.text = f"{value:.1f}"
        val.text_frame.paragraphs[0].font.size = Pt(9)
        val.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(DARK)


def add_range_chart(slide, x, y, w, h, ranges, min_axis, max_axis):
    row_h = h / len(ranges)
    axis_w = w - 2.2
    for i, (label, low, mid, high, color) in enumerate(ranges):
        yy = y + i * row_h
        lbl = slide.shapes.add_textbox(Inches(x), Inches(yy), Inches(1.9), Inches(row_h))
        lbl.text_frame.text = label
        lbl.text_frame.paragraphs[0].font.size = Pt(10)
        lbl.text_frame.paragraphs[0].font.bold = True
        lbl.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(DARK)
        def pos(value):
            return x + 2.0 + ((value - min_axis) / (max_axis - min_axis)) * axis_w
        start = pos(low)
        end = pos(high)
        bar = slide.shapes.add_shape(1, Inches(start), Inches(yy + 0.18), Inches(max(0.05, end - start)), Inches(0.18))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(color)
        bar.line.fill.background()
        dot = slide.shapes.add_shape(9, Inches(pos(mid) - 0.07), Inches(yy + 0.11), Inches(0.14), Inches(0.32))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor.from_string(DARK)
        dot.line.fill.background()
        txt = slide.shapes.add_textbox(Inches(end + 0.08), Inches(yy + 0.02), Inches(1.0), Inches(row_h))
        txt.text_frame.text = f"{low:.1f} - {high:.1f}"
        txt.text_frame.paragraphs[0].font.size = Pt(8.5)
        txt.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(GRAY)
    for tick in [30, 35, 40, 45, 50, 55, 60]:
        if min_axis <= tick <= max_axis:
            xx = x + 2.0 + ((tick - min_axis) / (max_axis - min_axis)) * axis_w
            line = slide.shapes.add_shape(1, Inches(xx), Inches(y + h + 0.08), Inches(0.01), Inches(0.10))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor.from_string(GRAY)
            line.line.fill.background()
            tick_txt = slide.shapes.add_textbox(Inches(xx - 0.18), Inches(y + h + 0.20), Inches(0.5), Inches(0.2))
            tick_txt.text_frame.text = str(tick)
            tick_txt.text_frame.paragraphs[0].font.size = Pt(7.5)
            tick_txt.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(GRAY)


def write_powerpoint(calcs: dict):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    hist = calcs["history"]["rows"]
    market = calcs["market"]
    wacc = calcs["wacc"]
    dcf = calcs["dcf"]
    rec = calcs["recommendation"]
    forecast = calcs["forecast"]

    # 1
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Michelin - Synthese d'investissement", "Valorisation par comparables et DCF, donnees au 20 avril 2026")
    add_metric(slide, 0.55, 1.15, 2.4, 1.0, eur_bn(hist["sales"][2025]), "CA 2025")
    add_metric(slide, 3.15, 1.15, 2.4, 1.0, pct(hist["segment_margin"][2025]), "Marge EBIT secteurs")
    add_metric(slide, 5.75, 1.15, 2.4, 1.0, eur_bn(hist["fcf"][2025]), "FCF 2025")
    add_metric(slide, 8.35, 1.15, 2.4, 1.0, eur(dcf["price"]), "Valeur DCF/action")
    add_metric(slide, 10.95, 1.15, 1.8, 1.0, rec["rating"], "Recommandation", RED)
    add_bullets(slide, 0.65, 2.55, 6.0, 3.5, [
        "Michelin reste un leader premium mondial avec une forte generation de cash-flow malgre la faiblesse des volumes OE en 2025.",
        "Le marche pneus est cyclique en premiere monte, mais le remplacement reste plus resilient et soutient le mix prix.",
        f"Objectif central retenu: {eur(rec['blended_target'])}/action, soit {pct(rec['blended_upside'])} de potentiel hors dividende.",
        f"Rendement total estime: {pct(rec['total_return'])}, superieur au cout des fonds propres de {pct(wacc['cost_equity'])}.",
    ], 14)
    add_bullets(slide, 7.1, 2.55, 5.4, 3.5, [
        "Principaux catalyseurs: reprise progressive des volumes, mix premium 18 pouces et plus, specialites, rachats d'actions 2026-2028.",
        "Risques majeurs: pression des importations a bas prix, couts matieres, change, faiblesse durable en premiere monte et sous-utilisation industrielle.",
        "Conclusion: Acheter, avec une marge de securite correcte face au cours actuel.",
    ], 14)
    add_footer(slide, "Sources: Michelin 2025 Annual Results, key figures 2020-2025, StockAnalysis, Kroll, CountryEconomy.")

    # 2
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Marche adressable et segmentation", "Un marche mondial important, tire par le remplacement et le premium")
    add_metric(slide, 0.6, 1.05, 2.7, 0.95, "USD 147,4 Md", "Marche pneus 2025")
    add_metric(slide, 3.55, 1.05, 2.7, 0.95, "USD 173,9 Md", "Prevision 2030")
    add_metric(slide, 6.5, 1.05, 2.7, 0.95, "3,4%", "CAGR 2025-2030")
    add_metric(slide, 9.45, 1.05, 2.7, 0.95, "Remplacement", "Canal le plus resilient", BLUE)
    markets = calcs["history"]["markets"]
    market_2025 = [
        ("PLT remplacement", markets["PLT Replacement"][2025]),
        ("PLT OE", markets["PLT OE"][2025]),
        ("Truck remplacement", markets["Truck Replacement"][2025]),
        ("Truck OE", markets["Truck OE"][2025]),
    ]
    add_horizontal_bars(slide, 0.7, 2.45, 6.0, 2.6, market_2025, max(v for _, v in market_2025))
    add_table(slide, 7.2, 2.2, 5.4, 2.4, [
        ["Segment", "Logique economique"],
        ["Premiere monte", "Depend de la production auto/camion"],
        ["Remplacement", "Parc roulant, kilometrage, securite"],
        ["Specialites", "Mining, aviation, agricole, deux-roues"],
        ["Non-pneus", "Connected fleet, polymeres, lifestyle"],
    ], [1.55, 3.75], 9)
    add_footer(slide, "Sources: Grand View Research; Michelin key figures 2020-2025, onglet Markets.")

    # 3
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Concurrence et barrieres a l'entree", "Un secteur mondial concentre sur les marques premium, avec des challengers asiatiques")
    add_table(slide, 0.55, 1.15, 6.1, 3.4, [
        ["Concurrent", "Positionnement"],
        ["Bridgestone", "Leader global, pneus et solutions"],
        ["Michelin", "Premium, innovation, specialites"],
        ["Goodyear", "Pneus purs, exposition US forte"],
        ["Continental", "Equipementier auto diversifie"],
        ["Pirelli", "Premium consumer et performance"],
        ["Yokohama / Hankook", "Acteurs asiatiques en montee"],
    ], [2.0, 4.1], 9)
    add_bullets(slide, 7.05, 1.2, 5.5, 3.7, [
        "Barrieres a l'entree: marque, reseau de distribution, relations OEM, certification, R&D et capital industriel.",
        "La performance securite, bruit, efficience energetique et durabilite impose des investissements continus.",
        "Le risque concurrentiel vient surtout des pneus budget importes et de la pression prix en remplacement.",
        "Michelin defend ses marges via mix premium, innovation produit, specialites et discipline prix.",
    ], 14)
    add_footer(slide, "Sources: Michelin URD 2025; StockAnalysis comparables; analyse interne.")

    # 4
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Croissance historique, drivers et risques", "2025 faible en volumes Michelin, mais marche global plus stable en remplacement")
    add_table(slide, 0.55, 1.15, 6.1, 2.6, [
        ["Marche 2025", "Croissance"],
        ["PLT premiere monte global", "+2%"],
        ["PLT remplacement global", "+1%"],
        ["Truck OE global", "+6%"],
        ["Truck remplacement global", "+3%"],
        ["Guidance marche 2026 Michelin", "PLT/TBR -2% a +2%"],
    ], [3.4, 2.7], 10)
    add_bullets(slide, 7.05, 1.1, 5.6, 4.7, [
        "Drivers: premiumisation, pneus 18 pouces et plus, electrification, efficience energetique, flottes connectees et specialites.",
        "Le remplacement beneficie du parc roulant et du kilometrage, moins volatil que la premiere monte.",
        "Risques: faiblesse auto/camion, matieres premieres, change EUR/USD, importations asiatiques, droits de douane et sous-utilisation des usines.",
        "Hypothese centrale: reprise graduelle, sans acceleration macro agressive.",
    ], 14)
    add_footer(slide, "Sources: Michelin 2025 Annual Results; Michelin market trends; European Rubber Journal.")

    # 5
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Activite, actionnariat et management", "Michelin combine pneus, composites, connected solutions et lifestyle")
    add_metric(slide, 0.6, 1.05, 2.4, 0.85, "115 800", "Employes FTE")
    add_metric(slide, 3.2, 1.05, 2.4, 0.85, "75", "Sites pneus")
    add_metric(slide, 5.8, 1.05, 2.4, 0.85, "48,8%", "Institutions")
    add_metric(slide, 8.4, 1.05, 2.4, 0.85, ">250k", "Actionnaires individuels")
    add_table(slide, 0.65, 2.25, 5.9, 2.6, [
        ["Segment 2025", "EBIT secteurs", "Marge"],
        ["Auto et distribution", eur_bn(calcs["history"]["segments"]["Automobile et distribution associee"][2025]["segment_ebit"]), "11,7%"],
        ["Transport routier", eur_bn(calcs["history"]["segments"]["Transport routier et distribution associee"][2025]["segment_ebit"]), "4,7%"],
        ["Specialites", eur_bn(calcs["history"]["segments"]["Specialites et distribution associee"][2025]["segment_ebit"]), "13,5%"],
    ], [2.5, 1.7, 1.2], 9.5)
    add_bullets(slide, 7.05, 2.2, 5.5, 2.9, [
        "Direction actuelle: Florent Menegaux, President de la Gerance et Associe Commandite; Yves Chapot, Gerant non Commandite.",
        "Philippe Jacquin est propose comme futur Gerant non Commandite a l'AG du 22 mai 2026.",
        "Strategie: valeur plutot que volume, leadership de marque MICHELIN, specialites et polymer composite solutions.",
    ], 13)
    add_footer(slide, "Sources: Michelin governance pages, key figures workbook, StockAnalysis.")

    # 6
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Analyse du P&L historique", "Recul des ventes 2024-2025, mais cash-flow et bilan solides")
    add_table(slide, 0.5, 1.05, 7.2, 3.2, [
        ["EURm", "2020", "2021", "2022", "2023", "2024", "2025"],
        ["CA", f"{hist['sales'][2020]:,.0f}", f"{hist['sales'][2021]:,.0f}", f"{hist['sales'][2022]:,.0f}", f"{hist['sales'][2023]:,.0f}", f"{hist['sales'][2024]:,.0f}", f"{hist['sales'][2025]:,.0f}"],
        ["Croissance", pct(hist["sales_growth"][2020]), pct(hist["sales_growth"][2021]), pct(hist["sales_growth"][2022]), pct(hist["sales_growth"][2023]), pct(hist["sales_growth"][2024]), pct(hist["sales_growth"][2025])],
        ["Marge EBIT secteurs", pct(hist["segment_margin"][2020]), pct(hist["segment_margin"][2021]), pct(hist["segment_margin"][2022]), pct(hist["segment_margin"][2023]), pct(hist["segment_margin"][2024]), pct(hist["segment_margin"][2025])],
        ["Resultat net", f"{hist['net_income'][2020]:,.0f}", f"{hist['net_income'][2021]:,.0f}", f"{hist['net_income'][2022]:,.0f}", f"{hist['net_income'][2023]:,.0f}", f"{hist['net_income'][2024]:,.0f}", f"{hist['net_income'][2025]:,.0f}"],
        ["FCF", f"{hist['fcf'][2020]:,.0f}", f"{hist['fcf'][2021]:,.0f}", f"{hist['fcf'][2022]:,.0f}", f"{hist['fcf'][2023]:,.0f}", f"{hist['fcf'][2024]:,.0f}", f"{hist['fcf'][2025]:,.0f}"],
    ], [1.4, 0.9, 0.9, 0.9, 0.9, 0.9, 0.9], 8)
    add_bullets(slide, 8.05, 1.1, 4.7, 4.0, [
        "2025: CA EUR 26,0 Md, en baisse de 4,4% a change courant.",
        "Volumes pneus -4,7%, principalement en premiere monte camion et agricole en Amerique du Nord.",
        "Marge secteurs 10,5%, penalisee par la sous-utilisation industrielle.",
        "FCF EUR 2,2 Md et dette nette reduite a EUR 2,3 Md.",
    ], 13)
    add_footer(slide, "Source: Michelin key figures 2020-2025.")

    # 7
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Projection du P&L sur 5 ans", "La feuille Forecast du classeur Excel sert de reference de verite")
    add_table(slide, 0.5, 1.05, 7.5, 3.4, [
        ["EURm", "2026E", "2027E", "2028E", "2029E", "2030E"],
        ["CA", *[f"{forecast[y]['sales']:,.0f}" for y in [2026, 2027, 2028, 2029, 2030]]],
        ["Croissance", *[pct(forecast[y]["growth"]) for y in [2026, 2027, 2028, 2029, 2030]]],
        ["Marge EBIT SECTOR MICHELIN", *[pct(forecast[y]["ebit_margin"]) for y in [2026, 2027, 2028, 2029, 2030]]],
        ["EBIT SECTOR MICHELIN", *[f"{forecast[y]['ebit']:,.0f}" for y in [2026, 2027, 2028, 2029, 2030]]],
        ["FCFF", *[f"{forecast[y]['fcff']:,.0f}" for y in [2026, 2027, 2028, 2029, 2030]]],
    ], [1.5, 1.2, 1.2, 1.2, 1.2, 1.2], 8.5)
    add_bullets(slide, 8.35, 1.1, 4.4, 4.1, [
        "2026E: premier point de projection du classeur Excel, repris tel quel dans la slide.",
        "Le CA progresse de facon prudente et la marge EBIT SECTOR MICHELIN remonte graduellement jusqu'a 12,0% en 2030.",
        "Le FCFF reste soutenu car le capex est maintenu proche de 7% du CA et la variation de BFR reste limitee.",
        "Les hypotheses de D&A / CA et de Capex / CA suivent la trajectoire incluse dans la feuille Forecast.",
    ], 12.5)
    add_footer(slide, "Source: feuille Forecast du classeur Excel Michelin.")

    # 8
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Echantillon de comparables retenu", "Pairs cotes mondiaux de pneumatiques et equipement automobile")
    add_table(slide, 0.4, 0.95, 12.5, 2.15, [
        ["Societe", "Activite", "Pourquoi retenu", "Limite"],
        *[[c["company"], c["activity"], c["relevance"], c["limitation"]] for c in calcs["comps"]],
    ], [1.6, 2.5, 4.25, 4.15], 6.9)
    add_table(slide, 0.4, 3.45, 12.5, 1.9, [
        ["Societe", "EV / EqV", "EV / CA", "EV / EBITDA", "EV / EBIT", "Marge EBIT"],
        *[[c["company"], f"{c['ev_eqv']:.2f}x", f"{c['ev_sales']:.2f}x", f"{c['ev_ebitda']:.2f}x", f"{c['ev_ebit']:.2f}x", pct(c["ebit_margin"])] for c in calcs["comps"]],
    ], [2.3, 1.3, 1.45, 1.65, 1.45, 1.45], 7.6)
    add_bullets(slide, 0.55, 5.55, 12.0, 0.88, [
        "Bridgestone et Pirelli sont les comparables les plus propres en termes d'activite et de qualite de marge.",
        "Goodyear reste utile sectoriellement, mais son EV / EqV eleve et sa faible marge rendent son EV / EBIT peu representatif.",
    ], 10.9)
    add_footer(slide, "Sources: StockAnalysis, statistiques de valorisation des comparables au 20 avril 2026 ou a la derniere cloture disponible.")

    # 9
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Resultats des comparables", "La mediane des pairs implique une valeur proche de EUR 36/action")
    tm = calcs["michelin_trading_multiples"]
    comp_rows = [
        ["Multiple", "Quartile bas", "Mediane", "Quartile haut", "Michelin actuel"],
        ["EV / CA", f"{calcs['comp_metrics']['EV / CA']['low']:.2f}x", f"{calcs['comp_metrics']['EV / CA']['mid']:.2f}x", f"{calcs['comp_metrics']['EV / CA']['high']:.2f}x", f"{tm['ev_sales']:.2f}x"],
        ["EV / EBITDA", f"{calcs['comp_metrics']['EV / EBITDA']['low']:.2f}x", f"{calcs['comp_metrics']['EV / EBITDA']['mid']:.2f}x", f"{calcs['comp_metrics']['EV / EBITDA']['high']:.2f}x", f"{tm['ev_ebitda']:.2f}x"],
        ["EV / EBIT", f"{calcs['comp_metrics']['EV / EBIT']['low']:.2f}x", f"{calcs['comp_metrics']['EV / EBIT']['mid']:.2f}x", f"{calcs['comp_metrics']['EV / EBIT']['high']:.2f}x", f"{tm['ev_ebit']:.2f}x"],
        ["EV / EqV", f"{calcs['peer_stats']['ev_eqv_low']:.2f}x", f"{calcs['peer_stats']['median_ev_eqv']:.2f}x", f"{calcs['peer_stats']['ev_eqv_high']:.2f}x", f"{tm['ev_eqv']:.2f}x"],
    ]
    add_table(slide, 0.65, 1.15, 6.6, 2.55, comp_rows, [1.8, 1.2, 1.2, 1.2, 1.2], 9.0)
    add_metric(slide, 7.7, 1.2, 2.15, 0.95, eur(rec["comps_mid_price"]), "Prix median comparables")
    add_metric(slide, 10.0, 1.2, 2.15, 0.95, pct(rec["comps_mid_price"] / market["share_price"] - 1), "Potentiel vs cours")
    ranges = []
    for method, item in calcs["comp_metrics"].items():
        ranges.append((method, item["price_low"], item["price_mid"], item["price_high"], TEAL))
    add_range_chart(slide, 0.75, 4.15, 11.2, 1.45, ranges, 25, 45)
    add_bullets(slide, 7.45, 3.55, 5.0, 1.55, [
        f"Michelin traite a {tm['ev_sales']:.2f}x EV / CA contre {calcs['comp_metrics']['EV / CA']['mid']:.2f}x pour la mediane des pairs.",
        f"Le EV / EBITDA Michelin utilise un EBITDA 2026E estime en interne: EBIT secteurs consensus ({eur(tm['segment_ebit_2026_consensus'])}) + D&A estimes.",
        f"Les D&A 2026E sont estimes comme {pct(tm['da_margin_2026_assumed'])} du CA 2026E consensus ({eur_bn(tm['sales_2026_consensus'])}), soit un taux proche de l'historique recent.",
        f"Sur EV / EBIT, Michelin est a {tm['ev_ebit']:.2f}x, proche du quartile bas ({calcs['comp_metrics']['EV / EBIT']['low']:.2f}x).",
        f"Michelin est a {tm['ev_eqv']:.2f}x EV / EqV contre {calcs['peer_stats']['median_ev_eqv']:.2f}x pour la mediane des pairs; ce ratio sert a lire le levier, pas a produire un prix implicite.",
    ], 9.6)
    add_footer(slide, "Sources: StockAnalysis peer multiples; Michelin 2026E consensus sales and segment EBIT; EBITDA 2026E Michelin estime en interne a partir de l'EBIT secteurs consensus et d'un ratio D&A / CA de 7,3% proche de l'historique recent.")

    # 10
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "WACC et hypotheses DCF", "WACC central de 8,8% et croissance terminale de 1,5%")
    add_table(slide, 0.6, 1.05, 5.9, 3.6, [
        ["Input", "Valeur"],
        ["Taux sans risque France 10 ans", pct(wacc["risk_free_rate"])],
        ["Prime de risque actions Eurozone", pct(wacc["erp"])],
        ["Beta Michelin", f"{wacc['beta']:.2f}x"],
        ["Cout des fonds propres", pct(wacc["cost_equity"])],
        ["Cout dette apres impots", pct(wacc["after_tax_cost_debt"])],
        ["Poids equity / dette nette", f"{pct(wacc['weight_equity'])} / {pct(wacc['weight_debt'])}"],
        ["WACC retenu", pct(wacc["wacc"])],
        ["Croissance terminale", pct(wacc["terminal_growth"])],
    ], [3.8, 1.8], 9.5)
    add_bullets(slide, 7.05, 1.25, 5.4, 3.7, [
        "Le DCF utilise l'EBIT secteurs comme EBIT normalise, afin de neutraliser les charges non recurrentes.",
        "Taux d'impot 26,3%, egal au taux effectif 2025.",
        "Capex converge vers 7,0% du CA, proche de l'historique Michelin.",
        "Croissance terminale 1,5%, prudente face a un marche pneus attendu autour de 3,4% CAGR nominal mondial.",
    ], 13)
    add_footer(slide, "Sources: CountryEconomy, Kroll, Michelin, StockAnalysis.")

    # 11
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Resultats du DCF", "La valeur terminale reste le principal contributeur, comme attendu pour un actif industriel mature")
    add_metric(slide, 0.6, 1.1, 2.35, 0.9, eur_bn(dcf["pv_fcff"]), "PV FCFF explicites")
    add_metric(slide, 3.25, 1.1, 2.35, 0.9, eur_bn(dcf["pv_terminal"]), "PV valeur terminale")
    add_metric(slide, 5.9, 1.1, 2.35, 0.9, eur_bn(dcf["enterprise_value"]), "EV DCF")
    add_metric(slide, 8.55, 1.1, 2.35, 0.9, eur(dcf["price"]), "Valeur/action")
    add_metric(slide, 11.2, 1.1, 1.55, 0.9, pct(dcf["upside"]), "Upside", RED)
    add_table(slide, 0.65, 2.55, 6.2, 2.5, [
        ["EURm", "2026E", "2027E", "2028E", "2029E", "2030E"],
        ["FCFF", *[f"{dcf['detail'][y]['fcff']:,.0f}" for y in dcf["years"]]],
        ["PV FCFF", *[f"{dcf['detail'][y]['pv_fcff']:,.0f}" for y in dcf["years"]]],
    ], [1.2, 1.0, 1.0, 1.0, 1.0, 1.0], 8.5)
    sens_mid = dcf["sensitivity"][wacc["wacc"]][wacc["terminal_growth"]]
    add_bullets(slide, 7.35, 2.65, 5.2, 2.4, [
        f"Valeur terminale brute: {eur_bn(dcf['terminal_value'])}.",
        f"Equity value DCF: {eur_bn(dcf['equity_value'])}.",
        f"Sensibilite centrale WACC/g retrouve {eur(sens_mid)} par action.",
        "Le resultat est sensible au WACC, d'ou la comparaison avec les comparables.",
    ], 12.5)
    add_footer(slide, "Source: modele Excel DCF.")

    # 12
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Football field et recommandation", "Recommendation: Acheter")
    comp_low = min(item["price_low"] for item in calcs["comp_metrics"].values())
    comp_high = max(item["price_high"] for item in calcs["comp_metrics"].values())
    dcf_prices = [p for row in dcf["sensitivity"].values() for p in row.values()]
    ranges = [
        ("Cours actuel", market["share_price"], market["share_price"], market["share_price"], GRAY),
        ("Comparables", comp_low, rec["comps_mid_price"], comp_high, TEAL),
        ("DCF", min(dcf_prices), dcf["price"], max(dcf_prices), BLUE),
        ("Objectif", rec["blended_target"], rec["blended_target"], rec["blended_target"], RED),
    ]
    add_range_chart(slide, 0.75, 1.25, 10.9, 2.4, ranges, 25, 60)
    add_metric(slide, 0.75, 4.45, 2.5, 0.9, eur(rec["blended_target"]), "Objectif central")
    add_metric(slide, 3.55, 4.45, 2.5, 0.9, pct(rec["blended_upside"]), "Potentiel cours")
    add_metric(slide, 6.35, 4.45, 2.5, 0.9, pct(market["dividend_yield"]), "Rendement dividende")
    add_metric(slide, 9.15, 4.45, 2.5, 0.9, pct(rec["total_return"]), "Rendement total", RED)
    add_bullets(slide, 0.85, 5.75, 11.7, 0.9, [
        "Acheter: le rendement total attendu depasse nettement le cout des fonds propres, avec un bilan solide et un FCF eleve.",
        "A surveiller: volumes OE, pression des pneus budget, change, matieres premieres et execution des rachats d'actions.",
    ], 11.5)
    add_footer(slide, "Sources: modele Excel, StockAnalysis, Michelin.")

    prs.save(PPT_OUT)


def write_summary(calcs: dict):
    lines = [
        "# Synthese de valorisation Michelin",
        "",
        f"- Cours de reference: {eur(calcs['market']['share_price'])}/action.",
        f"- Valeur par action comparables: {eur(calcs['recommendation']['comps_mid_price'])}.",
        f"- Valeur par action DCF: {eur(calcs['dcf']['price'])}.",
        f"- Objectif central retenu: {eur(calcs['recommendation']['blended_target'])}.",
        f"- Potentiel hors dividende: {pct(calcs['recommendation']['blended_upside'])}.",
        f"- Rendement total avec dividende: {pct(calcs['recommendation']['total_return'])}.",
        f"- Recommandation: {calcs['recommendation']['rating']}.",
        "",
        "Livrables generes:",
        "",
        f"- `{EXCEL_OUT.relative_to(ROOT)}`",
        f"- `{PPT_OUT.relative_to(ROOT)}`",
    ]
    SUMMARY_OUT.write_text("\n".join(lines) + "\n", encoding="utf-8")


def main():
    for path in [SOURCES, EXCEL_OUT.parent, PPT_OUT.parent, SUMMARY_OUT.parent]:
        path.mkdir(parents=True, exist_ok=True)
    history = read_michelin_history()
    calcs = build_calculations(history)
    write_excel(calcs)
    write_powerpoint(calcs)
    write_summary(calcs)
    print(f"Excel: {EXCEL_OUT}")
    print(f"PowerPoint: {PPT_OUT}")
    print(f"Summary: {SUMMARY_OUT}")
    print(f"DCF price: {calcs['dcf']['price']:.2f}")
    print(f"Blended target: {calcs['recommendation']['blended_target']:.2f}")
    print(f"Slides: 12")


if __name__ == "__main__":
    main()
